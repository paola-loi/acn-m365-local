from __future__ import annotations

import html
import os
import queue
import sys
import threading
import time
from datetime import datetime, timedelta
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from pathlib import Path
from typing import Optional
from urllib.parse import parse_qs, quote, unquote, urlsplit

from mcp.server.fastmcp import FastMCP
from mcp.server.fastmcp.exceptions import ToolError

SUPPORTED_FORMATS = [".docx", ".xlsx", ".pptx", ".md", ".txt"]


# ── path helpers ───────────────────────────────────────────────────

def _get_onedrive_base() -> Path:
    username = os.environ.get("USERNAME") or os.environ.get("USER", "")
    base = Path(f"C:/Users/{username}/OneDrive - Accenture")
    if not base.exists():
        raise ToolError(
            f"OneDrive base path not found: {base}. "
            "Verify that OneDrive is synced and the USERNAME environment variable is set."
        )
    return base


def _resolve_path(relative: str) -> Path:
    base = _get_onedrive_base()
    target = (base / relative).resolve()
    if not str(target).startswith(str(base.resolve())):
        raise ToolError(f"Path {relative!r} escapes the OneDrive directory.")
    return target


# ── text extraction ────────────────────────────────────────────────

def _read_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _read_xlsx(path: Path) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(str(path), read_only=True, data_only=True)
    lines: list[str] = []
    for sheet in wb.worksheets:
        lines.append(f"## Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                lines.append("\t".join(cells))
    wb.close()
    return "\n".join(lines)


def _xlsx_workbook(path: Path, *, data_only: bool, read_only: bool = True):
    from openpyxl import load_workbook
    return load_workbook(str(path), read_only=read_only, data_only=data_only)


def _read_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)


def _read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _extract_text(path: Path) -> str:
    ext = path.suffix.lower()
    try:
        match ext:
            case ".docx":
                return _read_docx(path)
            case ".xlsx":
                return _read_xlsx(path)
            case ".pptx":
                return _read_pptx(path)
            case ".md" | ".txt":
                return _read_text(path)
            case _:
                raise ToolError(
                    f"Unsupported file type '{ext}'. "
                    f"Supported formats: {', '.join(SUPPORTED_FORMATS)}"
                )
    except PermissionError:
        raise ToolError(
            f"Permission denied reading {path}. "
            "Verify OneDrive sync is complete and the file is not locked by another application."
        )


# ── HTTP viewer ────────────────────────────────────────────────────

_viewer_port: Optional[int] = None
_viewer_subscribers: dict[str, list["queue.Queue[str]"]] = {}
_viewer_mtimes: dict[str, float] = {}
_viewer_lock = threading.Lock()

_VIEWER_MAX_ROWS = 1000
_VIEWER_MAX_COLS = 100


def _viewer_key(path: Path) -> str:
    return str(path.resolve()).lower()


def _notify_viewer(path: Path) -> None:
    key = _viewer_key(path)
    with _viewer_lock:
        try:
            _viewer_mtimes[key] = path.stat().st_mtime
        except OSError:
            pass
        subs = list(_viewer_subscribers.get(key, ()))
    for q in subs:
        try:
            q.put_nowait("reload")
        except queue.Full:
            pass


def _viewer_poll_loop() -> None:
    while True:
        time.sleep(1.0)
        with _viewer_lock:
            keys = list(_viewer_subscribers.keys())
        for key in keys:
            try:
                mtime = os.stat(key).st_mtime
            except OSError:
                continue
            prev = _viewer_mtimes.get(key)
            if prev is None:
                _viewer_mtimes[key] = mtime
                continue
            if mtime != prev:
                _viewer_mtimes[key] = mtime
                with _viewer_lock:
                    subs = list(_viewer_subscribers.get(key, ()))
                for q in subs:
                    try:
                        q.put_nowait("reload")
                    except queue.Full:
                        pass


def _render_xlsx_html(target: Path, relpath: str) -> str:
    wb_v = _xlsx_workbook(target, data_only=True, read_only=True)
    wb_f = _xlsx_workbook(target, data_only=False, read_only=True)
    try:
        sheet_blocks: list[str] = []
        for ws_v in wb_v.worksheets:
            ws_f = wb_f[ws_v.title]
            max_row = min(ws_v.max_row or 0, _VIEWER_MAX_ROWS)
            max_col = min(ws_v.max_column or 0, _VIEWER_MAX_COLS)
            truncated = (
                (ws_v.max_row or 0) > _VIEWER_MAX_ROWS
                or (ws_v.max_column or 0) > _VIEWER_MAX_COLS
            )

            f_rows_iter = ws_f.iter_rows(
                min_row=1, max_row=max_row, min_col=1, max_col=max_col
            )
            rows_html: list[str] = []
            for v_row in ws_v.iter_rows(
                min_row=1, max_row=max_row, min_col=1, max_col=max_col
            ):
                f_row = next(f_rows_iter, ())
                tds: list[str] = []
                for i, v_cell in enumerate(v_row):
                    v = v_cell.value
                    f = f_row[i].value if i < len(f_row) else None
                    if v is None and isinstance(f, str) and f.startswith("="):
                        display = f"[{f}]"
                        cls = "formula"
                    else:
                        display = "" if v is None else str(v)
                        cls = "formula" if isinstance(f, str) and f.startswith("=") else ""
                    tds.append(
                        f'<td class="{cls}" title="{html.escape(v_cell.coordinate)}">'
                        f"{html.escape(display)}</td>"
                    )
                rows_html.append("<tr>" + "".join(tds) + "</tr>")

            banner = (
                f'<div class="banner">Truncated to {max_row}×{max_col} '
                f"(actual {ws_v.max_row}×{ws_v.max_column})</div>"
                if truncated
                else ""
            )
            sheet_blocks.append(
                f'<section class="sheet" data-name="{html.escape(ws_v.title)}">'
                f"<h2>{html.escape(ws_v.title)}</h2>"
                f"{banner}"
                f'<div class="scroll"><table>{"".join(rows_html)}</table></div>'
                f"</section>"
            )
    finally:
        wb_v.close()
        wb_f.close()

    tabs_html = "".join(
        f'<button data-target="{html.escape(ws.title)}">{html.escape(ws.title)}</button>'
        for ws in wb_v.worksheets
    )

    events_url = f"/events?path={quote(relpath)}"
    return f"""<!doctype html>
<html><head><meta charset="utf-8"><title>{html.escape(relpath)}</title>
<style>
body{{font:13px/1.4 system-ui,sans-serif;margin:0;padding:0;background:#fafafa;color:#222}}
header{{padding:8px 16px;background:#fff;border-bottom:1px solid #ddd;position:sticky;top:0;z-index:10}}
header h1{{font-size:14px;margin:0 0 4px;font-weight:600}}
header .status{{font-size:11px;color:#888}}
nav{{padding:4px 16px;background:#fff;border-bottom:1px solid #eee}}
nav button{{margin-right:4px;padding:4px 10px;border:1px solid #ccc;background:#fff;cursor:pointer;border-radius:3px;font-size:12px}}
nav button.active{{background:#0066cc;color:#fff;border-color:#0066cc}}
.sheet{{display:none;padding:16px}}
.sheet.active{{display:block}}
.sheet h2{{font-size:13px;margin:0 0 8px;color:#555}}
.scroll{{overflow:auto;max-height:calc(100vh - 140px);background:#fff;border:1px solid #ddd}}
table{{border-collapse:collapse;font-size:12px}}
td{{border:1px solid #e0e0e0;padding:3px 6px;white-space:nowrap;max-width:300px;overflow:hidden;text-overflow:ellipsis}}
td.formula{{background:#fff8e1}}
.banner{{background:#fff3cd;border:1px solid #ffe69c;padding:4px 8px;margin-bottom:8px;font-size:11px;border-radius:3px}}
</style></head>
<body>
<header>
  <h1>{html.escape(relpath)}</h1>
  <div class="status" id="status">connecting…</div>
</header>
<nav>{tabs_html}</nav>
{"".join(sheet_blocks)}
<script>
(function(){{
  var tabs=document.querySelectorAll('nav button');
  var sheets=document.querySelectorAll('.sheet');
  function activate(name){{
    tabs.forEach(function(b){{b.classList.toggle('active',b.dataset.target===name)}});
    sheets.forEach(function(s){{s.classList.toggle('active',s.dataset.name===name)}});
  }}
  if(tabs.length){{activate(tabs[0].dataset.target);
    tabs.forEach(function(b){{b.addEventListener('click',function(){{activate(b.dataset.target)}})}});
  }}
  var status=document.getElementById('status');
  var es=new EventSource({events_url!r});
  es.onopen=function(){{status.textContent='live • '+new Date().toLocaleTimeString()}};
  es.onmessage=function(e){{if(e.data==='reload')location.reload();}};
  es.onerror=function(){{status.textContent='disconnected — retrying'}};
}})();
</script>
</body></html>"""


class _ViewerHandler(BaseHTTPRequestHandler):
    def log_message(self, fmt, *args):
        return

    def _send(self, code: int, body: bytes, ctype: str = "text/html; charset=utf-8") -> None:
        self.send_response(code)
        self.send_header("Content-Type", ctype)
        self.send_header("Content-Length", str(len(body)))
        self.send_header("Cache-Control", "no-store")
        self.end_headers()
        self.wfile.write(body)

    def do_GET(self):  # noqa: N802
        parts = urlsplit(self.path)
        path = parts.path

        if path == "/" or path == "/index":
            with _viewer_lock:
                keys = list(_viewer_subscribers.keys())
            items = "".join(f"<li>{html.escape(k)}</li>" for k in keys) or "<li><em>no files subscribed yet</em></li>"
            body = f"<!doctype html><meta charset=utf-8><title>viewer</title><h1>acn-m365 viewer</h1><ul>{items}</ul>".encode()
            self._send(200, body)
            return

        if path.startswith("/view/"):
            relpath = unquote(path[len("/view/"):])
            try:
                target = _resolve_path(relpath)
            except ToolError as e:
                self._send(400, str(e).encode(), "text/plain; charset=utf-8")
                return
            if not target.exists() or target.suffix.lower() != ".xlsx":
                self._send(404, b"Not found or not .xlsx", "text/plain; charset=utf-8")
                return
            try:
                body = _render_xlsx_html(target, relpath).encode("utf-8")
            except Exception as e:
                self._send(500, f"Render error: {e}".encode(), "text/plain; charset=utf-8")
                return
            try:
                _viewer_mtimes[_viewer_key(target)] = target.stat().st_mtime
            except OSError:
                pass
            self._send(200, body)
            return

        if path == "/events":
            qs = parse_qs(parts.query)
            relpath = unquote((qs.get("path") or [""])[0])
            try:
                target = _resolve_path(relpath)
            except ToolError:
                self._send(400, b"bad path", "text/plain")
                return
            key = _viewer_key(target)
            q: "queue.Queue[str]" = queue.Queue(maxsize=8)
            with _viewer_lock:
                _viewer_subscribers.setdefault(key, []).append(q)
                if key not in _viewer_mtimes:
                    try:
                        _viewer_mtimes[key] = target.stat().st_mtime
                    except OSError:
                        pass
            try:
                self.send_response(200)
                self.send_header("Content-Type", "text/event-stream")
                self.send_header("Cache-Control", "no-store")
                self.send_header("Connection", "keep-alive")
                self.end_headers()
                self.wfile.write(b": connected\n\n")
                self.wfile.flush()
                while True:
                    try:
                        msg = q.get(timeout=15)
                        self.wfile.write(f"data: {msg}\n\n".encode())
                        self.wfile.flush()
                    except queue.Empty:
                        self.wfile.write(b": ping\n\n")
                        self.wfile.flush()
            except (BrokenPipeError, ConnectionResetError):
                pass
            finally:
                with _viewer_lock:
                    if key in _viewer_subscribers:
                        try:
                            _viewer_subscribers[key].remove(q)
                        except ValueError:
                            pass
                        if not _viewer_subscribers[key]:
                            del _viewer_subscribers[key]
            return

        self._send(404, b"Not found", "text/plain; charset=utf-8")


def _start_viewer_server() -> None:
    global _viewer_port
    if _viewer_port is not None:
        return
    last_err: Optional[OSError] = None
    for port in list(range(8765, 8776)) + [0]:
        try:
            srv = ThreadingHTTPServer(("127.0.0.1", port), _ViewerHandler)
            _viewer_port = srv.server_address[1]
            threading.Thread(target=srv.serve_forever, daemon=True, name="viewer-http").start()
            threading.Thread(target=_viewer_poll_loop, daemon=True, name="viewer-poll").start()
            print(
                f"[viewer] listening on http://127.0.0.1:{_viewer_port}",
                file=sys.stderr,
                flush=True,
            )
            return
        except OSError as e:
            last_err = e
            continue
    print(f"[viewer] failed to start: {last_err}", file=sys.stderr, flush=True)


# ── server ─────────────────────────────────────────────────────────

mcp = FastMCP(
    "acn-m365-local",
    instructions=(
        "Reads and writes Microsoft 365 files from OneDrive-synced local folders. "
        "No authentication required — all operations are local file system access."
    ),
)


# ── READ tools ─────────────────────────────────────────────────────

@mcp.tool()
def list_project_files(project: str, file_type: Optional[str] = None) -> list[dict]:
    """List files inside an OneDrive project folder.

    Args:
        project: Subfolder name inside OneDrive (e.g. "MARS", "Unilever/2025")
        file_type: Optional extension filter — one of: .docx, .xlsx, .pptx, .md, .txt
    """
    folder = _resolve_path(project)
    if not folder.exists():
        raise ToolError(
            f"Project folder not found: {folder}. "
            "Check that the folder exists and OneDrive is fully synced."
        )
    if not folder.is_dir():
        raise ToolError(f"Path is not a directory: {folder}")

    results = []
    for f in sorted(folder.rglob("*")):
        if not f.is_file():
            continue
        if file_type and f.suffix.lower() != file_type.lower():
            continue
        stat = f.stat()
        results.append({
            "name": f.name,
            "path": str(f.relative_to(_get_onedrive_base())),
            "modified_date": datetime.fromtimestamp(stat.st_mtime).isoformat(),
            "size": stat.st_size,
        })
    return results


@mcp.tool()
def read_document(path: str) -> dict:
    """Read the full text content of a file from OneDrive.

    Args:
        path: File path relative to OneDrive base (e.g. "MARS/report.docx")
    """
    target = _resolve_path(path)
    if not target.exists():
        raise ToolError(
            f"File not found: {target}. "
            "Check the path and verify OneDrive is fully synced."
        )
    if not target.is_file():
        raise ToolError(f"Path is not a file: {target}")

    stat = target.stat()
    content = _extract_text(target)
    return {
        "path": path,
        "name": target.name,
        "extension": target.suffix.lower(),
        "size": stat.st_size,
        "modified_date": datetime.fromtimestamp(stat.st_mtime).isoformat(),
        "content": content,
        "char_count": len(content),
    }


@mcp.tool()
def search_documents(query: str, project: Optional[str] = None) -> list[dict]:
    """Search for a keyword across OneDrive file contents.

    Args:
        query: Search term (case-insensitive)
        project: Optional subfolder to restrict the search (strongly recommended for speed)
    """
    base = _get_onedrive_base()
    search_root = _resolve_path(project) if project else base

    if not search_root.exists():
        raise ToolError(f"Search path not found: {search_root}")

    query_lower = query.lower()
    results = []

    for f in search_root.rglob("*"):
        if not f.is_file() or f.suffix.lower() not in SUPPORTED_FORMATS:
            continue
        try:
            text = _extract_text(f)
        except ToolError:
            continue

        if query_lower not in text.lower():
            continue

        idx = text.lower().find(query_lower)
        start = max(0, idx - 100)
        end = min(len(text), idx + 200)
        excerpt = (
            ("..." if start > 0 else "")
            + text[start:end]
            + ("..." if end < len(text) else "")
        )

        results.append({
            "file": f.name,
            "path": str(f.relative_to(base)),
            "excerpt": excerpt,
            "relevance_score": text.lower().count(query_lower),
        })

    return sorted(results, key=lambda x: x["relevance_score"], reverse=True)


@mcp.tool()
def get_recent_files(days: int = 7, project: Optional[str] = None) -> list[dict]:
    """List supported files modified in the last N days.

    Args:
        days: Look-back window in days (default: 7)
        project: Optional subfolder to restrict the search
    """
    base = _get_onedrive_base()
    search_root = _resolve_path(project) if project else base
    cutoff = datetime.now() - timedelta(days=days)

    results = []
    for f in search_root.rglob("*"):
        if not f.is_file() or f.suffix.lower() not in SUPPORTED_FORMATS:
            continue
        stat = f.stat()
        mtime = datetime.fromtimestamp(stat.st_mtime)
        if mtime >= cutoff:
            results.append({
                "name": f.name,
                "path": str(f.relative_to(base)),
                "modified_date": mtime.isoformat(),
                "size": stat.st_size,
            })

    return sorted(results, key=lambda x: x["modified_date"], reverse=True)


# ── XLSX structured tools ──────────────────────────────────────────

def _require_xlsx(path: str) -> Path:
    target = _resolve_path(path)
    if not target.exists():
        raise ToolError(f"File not found: {target}.")
    if target.suffix.lower() != ".xlsx":
        raise ToolError(f"Not an xlsx file: {target}.")
    return target


@mcp.tool()
def list_sheets(path: str) -> list[dict]:
    """List worksheets in an xlsx file with their dimensions.

    Args:
        path: File path relative to OneDrive base (e.g. "MARS/budget.xlsx")
    """
    target = _require_xlsx(path)
    wb = _xlsx_workbook(target, data_only=True, read_only=True)
    try:
        return [
            {
                "name": ws.title,
                "max_row": ws.max_row,
                "max_col": ws.max_column,
                "dimensions": ws.dimensions,
            }
            for ws in wb.worksheets
        ]
    finally:
        wb.close()


@mcp.tool()
def describe_sheet(path: str, sheet: str, header_row: int = 1) -> dict:
    """Describe a worksheet: headers, inferred column types, non-empty row count, named ranges.

    Args:
        path: File path relative to OneDrive base
        sheet: Worksheet name
        header_row: Row number containing headers (default: 1)
    """
    target = _require_xlsx(path)
    wb = _xlsx_workbook(target, data_only=True, read_only=True)
    try:
        if sheet not in wb.sheetnames:
            raise ToolError(f"Sheet '{sheet}' not found. Available: {wb.sheetnames}")
        ws = wb[sheet]

        headers: list = []
        for row in ws.iter_rows(min_row=header_row, max_row=header_row, values_only=True):
            headers = list(row)
            break

        max_col = ws.max_column or 0
        samples: list[list] = [[] for _ in range(max_col)]
        non_empty_rows = 0
        sample_limit = 20
        for row in ws.iter_rows(min_row=header_row + 1, values_only=True):
            if any(c is not None and c != "" for c in row):
                non_empty_rows += 1
            for i, cell in enumerate(row[:max_col]):
                if cell is not None and cell != "" and len(samples[i]) < sample_limit:
                    samples[i].append(cell)

        def infer_type(values: list) -> str:
            if not values:
                return "empty"
            types = {type(v).__name__ for v in values}
            if types == {"int"} or types == {"float"} or types <= {"int", "float"}:
                return "number"
            if types == {"bool"}:
                return "bool"
            if types == {"datetime"} or types == {"date"}:
                return "date"
            if types == {"str"}:
                return "string"
            return "mixed:" + ",".join(sorted(types))

        columns = [
            {
                "header": headers[i] if i < len(headers) else None,
                "inferred_type": infer_type(samples[i]),
                "sample": samples[i][:5],
            }
            for i in range(max_col)
        ]

        named_ranges = []
        try:
            for dn in wb.defined_names:
                dest = wb.defined_names[dn]
                named_ranges.append({"name": dn, "value": str(dest.value)})
        except Exception:
            pass

        return {
            "sheet": sheet,
            "dimensions": ws.dimensions,
            "max_row": ws.max_row,
            "max_col": max_col,
            "non_empty_data_rows": non_empty_rows,
            "header_row": header_row,
            "columns": columns,
            "named_ranges": named_ranges,
        }
    finally:
        wb.close()


@mcp.tool()
def read_cells(path: str, sheet: str, range: str, include_formulas: bool = True) -> dict:
    """Read a cell range from an xlsx worksheet, returning values and formulas.

    Args:
        path: File path relative to OneDrive base
        sheet: Worksheet name
        range: A1-notation range (e.g. "A1:D50")
        include_formulas: If true, also return formula text alongside cached values
    """
    target = _require_xlsx(path)

    wb_values = _xlsx_workbook(target, data_only=True, read_only=True)
    wb_formulas = _xlsx_workbook(target, data_only=False, read_only=True) if include_formulas else None
    try:
        if sheet not in wb_values.sheetnames:
            raise ToolError(f"Sheet '{sheet}' not found. Available: {wb_values.sheetnames}")
        ws_v = wb_values[sheet]
        ws_f = wb_formulas[sheet] if wb_formulas is not None else None

        try:
            value_rows = list(ws_v[range])
        except (ValueError, KeyError) as e:
            raise ToolError(f"Invalid range {range!r}: {e}")
        formula_rows = list(ws_f[range]) if ws_f is not None else None

        cells: list[list[dict]] = []
        for r_idx, row in enumerate(value_rows):
            out_row: list[dict] = []
            for c_idx, cell in enumerate(row):
                entry: dict = {"addr": cell.coordinate, "value": cell.value}
                if formula_rows is not None:
                    f_cell = formula_rows[r_idx][c_idx]
                    fv = f_cell.value
                    if isinstance(fv, str) and fv.startswith("="):
                        entry["formula"] = fv
                out_row.append(entry)
            cells.append(out_row)

        return {"sheet": sheet, "range": range, "cells": cells}
    finally:
        wb_values.close()
        if wb_formulas is not None:
            wb_formulas.close()


@mcp.tool()
def write_cells(
    path: str,
    sheet: str,
    changes: list[dict],
    output_suffix: str = ".claude",
) -> dict:
    """Apply cell changes and save to a sibling copy (never overwrites the original).

    Args:
        path: Source xlsx path relative to OneDrive base
        sheet: Worksheet name to modify
        changes: List of {"cell": "B7", "value": ...} or {"cell": "B7", "formula": "=SUM(A1:A6)"}
        output_suffix: Suffix inserted before .xlsx for the output copy (default: ".claude")
    """
    target = _require_xlsx(path)

    out_path = target.with_name(f"{target.stem}{output_suffix}.xlsx")
    source = out_path if out_path.exists() else target

    wb = _xlsx_workbook(source, data_only=False, read_only=False)
    try:
        if sheet not in wb.sheetnames:
            raise ToolError(f"Sheet '{sheet}' not found. Available: {wb.sheetnames}")
        ws = wb[sheet]

        applied = 0
        for ch in changes:
            cell_addr = ch.get("cell")
            if not cell_addr:
                raise ToolError(f"Change missing 'cell' key: {ch}")
            if "formula" in ch:
                f = ch["formula"]
                if not isinstance(f, str) or not f.startswith("="):
                    raise ToolError(f"Formula for {cell_addr} must be a string starting with '=': {f!r}")
                ws[cell_addr] = f
            elif "value" in ch:
                ws[cell_addr] = ch["value"]
            else:
                raise ToolError(f"Change for {cell_addr} must include 'value' or 'formula'.")
            applied += 1

        try:
            wb.save(str(out_path))
        except PermissionError:
            raise ToolError(
                f"Permission denied writing {out_path}. Close the file in Excel and retry."
            )
    finally:
        wb.close()

    _notify_viewer(out_path)

    base = _get_onedrive_base()
    return {
        "output_path": str(out_path.relative_to(base)),
        "sheet": sheet,
        "changes_applied": applied,
    }


@mcp.tool()
def view_xlsx(path: str) -> dict:
    """Return a localhost URL to view an xlsx live in the browser.

    The page auto-reloads when Claude writes (via write_cells) and when the
    file is saved externally (LibreOffice/Excel). Bound to 127.0.0.1 only.

    Args:
        path: xlsx path relative to OneDrive base (e.g. "MARS/budget.claude.xlsx")
    """
    target = _require_xlsx(path)
    if _viewer_port is None:
        raise ToolError("Viewer HTTP server is not running. Check server.py startup logs.")
    return {
        "url": f"http://127.0.0.1:{_viewer_port}/view/{quote(path)}",
        "port": _viewer_port,
        "resolved": str(target),
    }


@mcp.tool()
def find_in_xlsx(path: str, query: str, sheet: Optional[str] = None) -> list[dict]:
    """Search a string (case-insensitive) across cell values and formulas of an xlsx.

    Args:
        path: File path relative to OneDrive base
        query: Substring to find
        sheet: Optional sheet name to restrict the search
    """
    target = _require_xlsx(path)
    q = query.lower()

    wb_v = _xlsx_workbook(target, data_only=True, read_only=True)
    wb_f = _xlsx_workbook(target, data_only=False, read_only=True)
    try:
        sheets = [sheet] if sheet else wb_v.sheetnames
        if sheet and sheet not in wb_v.sheetnames:
            raise ToolError(f"Sheet '{sheet}' not found. Available: {wb_v.sheetnames}")

        hits: list[dict] = []
        for s in sheets:
            ws_v = wb_v[s]
            ws_f = wb_f[s]
            f_iter = ws_f.iter_rows()
            for v_row in ws_v.iter_rows():
                f_row = next(f_iter, ())
                for i, v_cell in enumerate(v_row):
                    v = v_cell.value
                    f = f_row[i].value if i < len(f_row) else None
                    formula = f if isinstance(f, str) and f.startswith("=") else None
                    v_str = "" if v is None else str(v)
                    f_str = formula or ""
                    if q in v_str.lower() or q in f_str.lower():
                        hits.append({
                            "sheet": s,
                            "cell": v_cell.coordinate,
                            "value": v,
                            "formula": formula,
                        })
        return hits
    finally:
        wb_v.close()
        wb_f.close()


# ── WRITE tools ────────────────────────────────────────────────────

@mcp.tool()
def create_document(path: str, content: str, doc_format: str = "md") -> dict:
    """Create a new file in OneDrive. Fails if the file already exists.

    Args:
        path: Destination path relative to OneDrive base (e.g. "MARS/notes.md")
        content: Text content to write
        doc_format: "md" for plain Markdown or "docx" for Word document
    """
    target = _resolve_path(path)
    if target.exists():
        raise ToolError(
            f"File already exists: {target}. Use update_document to overwrite it."
        )

    target.parent.mkdir(parents=True, exist_ok=True)

    try:
        ext = target.suffix.lower()
        if doc_format == "docx" or ext == ".docx":
            from docx import Document
            doc = Document()
            for line in content.splitlines():
                doc.add_paragraph(line)
            doc.save(str(target))
        else:
            target.write_text(content, encoding="utf-8")
    except PermissionError:
        raise ToolError(
            f"Permission denied writing {target}. "
            "Verify OneDrive sync is complete and no other app has the file locked."
        )

    return {"created": str(target), "path": path, "format": doc_format}


@mcp.tool()
def update_document(path: str, content: str) -> dict:
    """Overwrite an existing file with new content. Fails if the file does not exist.

    Args:
        path: File path relative to OneDrive base
        content: New text content (replaces existing content entirely)
    """
    target = _resolve_path(path)
    if not target.exists():
        raise ToolError(
            f"File not found: {target}. Use create_document to create a new file."
        )

    ext = target.suffix.lower()
    try:
        if ext == ".docx":
            from docx import Document
            doc = Document()
            for line in content.splitlines():
                doc.add_paragraph(line)
            doc.save(str(target))
        elif ext in (".md", ".txt"):
            target.write_text(content, encoding="utf-8")
        else:
            raise ToolError(
                f"Cannot update '{ext}' files directly. "
                "Supported formats for update: .md, .txt, .docx"
            )
    except PermissionError:
        raise ToolError(
            f"Permission denied writing {target}. "
            "Verify the file is not open in another application."
        )

    return {"updated": str(target), "path": path}


@mcp.tool()
def create_summary_report(project: str, output_path: str) -> dict:
    """Generate a Markdown summary report aggregating all documents in a project folder.

    The report includes: file list, per-file summary (first 500 chars), and extracted
    action items (lines containing: todo, action, follow up, da fare, azione).

    Args:
        project: Source project folder name inside OneDrive
        output_path: Destination path for the .md report (relative to OneDrive)
    """
    base = _get_onedrive_base()
    project_dir = _resolve_path(project)

    if not project_dir.exists():
        raise ToolError(
            f"Project folder not found: {project_dir}. "
            "Check the folder name and OneDrive sync status."
        )

    ACTION_KEYWORDS = ["todo", "action", "follow up", "followup", "da fare", "azione"]

    files_info = []
    for f in sorted(project_dir.rglob("*")):
        if not f.is_file() or f.suffix.lower() not in SUPPORTED_FORMATS:
            continue
        try:
            text = _extract_text(f)
        except ToolError:
            text = "(could not read file)"

        action_items = [
            ln.strip()
            for ln in text.splitlines()
            if any(kw in ln.lower() for kw in ACTION_KEYWORDS)
        ]

        stat = f.stat()
        files_info.append({
            "name": f.name,
            "path": str(f.relative_to(base)),
            "modified": datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d"),
            "size": stat.st_size,
            "summary": text[:500].replace("\n", " ") + ("..." if len(text) > 500 else ""),
            "action_items": action_items[:10],
        })

    lines = [
        f"# Summary Report: {project}",
        f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}",
        "",
        f"## Files ({len(files_info)} documents)",
        "",
    ]
    for fi in files_info:
        lines += [
            f"### {fi['name']}",
            f"- **Path:** `{fi['path']}`",
            f"- **Modified:** {fi['modified']}",
            f"- **Size:** {fi['size']:,} bytes",
            "",
            f"**Summary:** {fi['summary']}",
            "",
        ]
        if fi["action_items"]:
            lines.append("**Action items:**")
            for item in fi["action_items"]:
                lines.append(f"- {item}")
            lines.append("")

    report_content = "\n".join(lines)

    out_target = _resolve_path(output_path)
    out_target.parent.mkdir(parents=True, exist_ok=True)
    try:
        out_target.write_text(report_content, encoding="utf-8")
    except PermissionError:
        raise ToolError(f"Permission denied writing report to {out_target}.")

    return {
        "report_path": output_path,
        "files_processed": len(files_info),
        "created": str(out_target),
    }


# ── entry point ────────────────────────────────────────────────────

_start_viewer_server()


if __name__ == "__main__":
    if "--http" in sys.argv:
        # HTTP/SSE mode: required by Accenture enterprise policy
        # (policy allows http://localhost* but blocks stdio custom servers)
        mcp.run(transport="sse")
    else:
        mcp.run()
